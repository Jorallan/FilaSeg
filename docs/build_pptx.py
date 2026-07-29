"""Build docs/pipeline_overview.pptx from the figures and schematics in
docs/figures/. Layout: a title slide, a stage-overview deck (block diagram +
the 4 stage in/out PNGs), and a per-stage deep dive that uses every schematic
PNG in docs/figures/schematics/. GIFs are skipped (PowerPoint plays them
poorly inline and the first frame doesn't tell the whole story).

Slide layout for image slides:
    +-------------------------------------+
    |  Title                              |
    |                                     |
    |   [image, centered, aspect-fit]     |
    |                                     |
    |  caption                            |
    +-------------------------------------+

Run:
    C:\\Repos\\venv_cnt\\Scripts\\python.exe docs\\build_pptx.py
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
FIG = DOCS / "figures"
SCHEMA = FIG / "schematics"
OUT = DOCS / "pipeline_overview.pptx"

# Bundle GIFs to embed (final ID, dominant stringart branch, rec_id)
BUNDLE_GIFS = [
    (1,  "branch 1 (0-15°)",   "rec_id 1"),
    (6,  "branch 3 (30-45°)",  "rec_id 6"),
    (7,  "branch 4 (45-60°)",  "rec_id 7"),
    (8,  "branch 4 (45-60°)",  "rec_id 8"),
    (10, "branch 3 (30-45°)",  "rec_id 10"),
]

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Visual reservations on every image slide
TITLE_BOX = (Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
CAPTION_BOX = (Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.7))
IMAGE_AREA = (Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.45))

TITLE_FONT = "Calibri"
TITLE_PT = 26
CAPTION_PT = 13
SECTION_PT = 36

ACCENT_RGB = RGBColor(0x19, 0x76, 0xD2)
DARK_RGB = RGBColor(0x21, 0x21, 0x21)
GREY_RGB = RGBColor(0x55, 0x55, 0x55)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_run(run, *, font=TITLE_FONT, pt=TITLE_PT, bold=False, color=DARK_RGB):
    run.font.name = font
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, *,
                 font=TITLE_FONT, pt=TITLE_PT, bold=False,
                 color=DARK_RGB, align="left"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = {"left": PP_ALIGN.LEFT,
                   "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    _set_run(run, font=font, pt=pt, bold=bold, color=color)
    return tx


def fit_image_in_box(prs, slide, image_path: Path, area):
    """Insert image scaled to fit inside `area` (left, top, w, h), centered."""
    left, top, area_w, area_h = area
    with Image.open(image_path) as im:
        iw, ih = im.size
    img_aspect = iw / ih
    box_aspect = area_w / area_h
    if img_aspect > box_aspect:
        new_w = area_w
        new_h = Emu(int(area_w * (ih / iw)))
    else:
        new_h = area_h
        new_w = Emu(int(area_h * (iw / ih)))
    x = Emu(int(left + (area_w - new_w) / 2))
    y = Emu(int(top + (area_h - new_h) / 2))
    slide.shapes.add_picture(str(image_path), x, y, width=new_w, height=new_h)


def make_blank_slide(prs) -> "Slide":
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs):
    s = make_blank_slide(prs)

    # Header bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_RGB
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text_box(s, Inches(0.5), Inches(0.55), Inches(12), Inches(0.9),
                 "FilaSeg", pt=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.7),
                 "Filament-quantification pipeline overview",
                 pt=24, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Block diagram below
    fit_image_in_box(prs, s, FIG / "01_block_diagram.png",
                     (Inches(0.4), Inches(2.4), Inches(12.5), Inches(4.0)))

    add_text_box(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
                 "Six pipeline boxes: input → stringart → preprocess → reconnect → postprocess → final outputs.",
                 pt=14, color=GREY_RGB, align="center")


def add_section_divider(prs, big_title: str, subtitle: str = ""):
    s = make_blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_RGB
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text_box(s, Inches(0.5), Inches(2.8), Inches(12), Inches(1.0),
                 big_title, pt=SECTION_PT, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text_box(s, Inches(0.5), Inches(3.8), Inches(12), Inches(1.0),
                     subtitle, pt=20, color=RGBColor(0xFF, 0xFF, 0xFF))


def add_image_slide(prs, title: str, image_path: Path, caption: str = ""):
    """Slide layout: title (movable text box), image (movable picture), caption
    (movable text box). All three are independent shapes you can click and
    drag in PowerPoint - no group, no background, no layout placeholder.
    """
    s = make_blank_slide(prs)
    # title
    add_text_box(s, *TITLE_BOX, title, pt=TITLE_PT, bold=True)
    # image
    fit_image_in_box(prs, s, image_path, IMAGE_AREA)
    # caption
    if caption:
        add_text_box(s, *CAPTION_BOX, caption,
                     pt=CAPTION_PT, color=GREY_RGB, align="center")


# HSV colors matching matplotlib's hsv colormap used in the schematic
def _hsv_colors(n: int) -> list[RGBColor]:
    import colorsys
    out = []
    for i in range(n):
        r, g, b = colorsys.hsv_to_rgb(i / n, 1.0, 1.0)
        out.append(RGBColor(int(r * 255), int(g * 255), int(b * 255)))
    return out


def add_native_angle_binning_slide(prs):
    """Native-shape version of the angle-binning slide.

    Every bin cell is a separate PowerPoint rectangle with its own text - you
    can click and move each one independently, change colors in the PPT UI,
    re-order them, etc. The example line and arrow are also separate shapes.
    """
    s = make_blank_slide(prs)
    add_text_box(s, *TITLE_BOX,
                 "Stage 1 — angle binning (12 bins, 15° wide)",
                 pt=TITLE_PT, bold=True)

    colors = _hsv_colors(12)

    # ── Fan of 12 lines, drawn as separate native lines ──────────────────
    import math
    cx = Inches(6.667)
    cy = Inches(3.5)
    L = Inches(1.6)
    for i in range(12):
        mid_deg = i * 15 + 7.5
        rad = math.radians(mid_deg)
        # endpoint coords (note: in PPT, +y is downward, so we invert)
        dx = Emu(int(L * math.cos(rad)))
        dy = Emu(int(L * math.sin(rad)))
        x1 = Emu(int(cx + dx))
        y1 = Emu(int(cy - dy))
        line = s.shapes.add_connector(1, cx, cy, x1, y1)  # 1 = straight
        line.line.color.rgb = colors[i]
        line.line.width = Pt(5)
        # bin badge at the tip
        badge_size = Inches(0.32)
        bx = Emu(int(x1 - badge_size / 2))
        by = Emu(int(y1 - badge_size / 2))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        badge.line.color.rgb = colors[i]
        badge.line.width = Pt(2)
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        _set_run(run, pt=12, bold=True, color=colors[i])

    # Baseline under the fan
    base = s.shapes.add_connector(1, Inches(4.3), Inches(3.5),
                                   Inches(9.0), Inches(3.5))
    base.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    base.line.width = Pt(1.5)
    add_text_box(s, Inches(8.95), Inches(3.55), Inches(0.6), Inches(0.35),
                 "0°", pt=11, color=GREY_RGB, align="left")
    add_text_box(s, Inches(6.55), Inches(3.55), Inches(0.3), Inches(0.35),
                 "90°", pt=11, color=GREY_RGB, align="center")
    add_text_box(s, Inches(3.9), Inches(3.55), Inches(0.6), Inches(0.35),
                 "180°", pt=11, color=GREY_RGB, align="right")

    # ── Horizontal legend bar: 12 individual cells, each a separate shape ──
    legend_top = Inches(5.2)
    legend_h = Inches(0.95)
    legend_x = Inches(0.6)
    cell_w = Inches((13.333 - 1.2) / 12)
    for i in range(12):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Emu(int(legend_x + i * cell_w)),
                                  legend_top, cell_w, legend_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors[i]
        cell.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.line.width = Pt(1.5)
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"bin {i+1}"
        _set_run(run, pt=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"{i*15}° – {(i+1)*15}°"
        _set_run(run2, pt=9, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Caption
    add_text_box(s, *CAPTION_BOX,
                 "Native shapes — click any bin cell or the fan badge and drag it. "
                 "Every element on this slide is an individually movable PowerPoint shape.",
                 pt=CAPTION_PT, color=GREY_RGB, align="center")


def add_gif_slide(prs, title: str, gif_path: Path, caption: str = ""):
    """A GIF slide. python-pptx inserts the .gif file as a picture; PowerPoint
    will animate it during slideshow mode (F5)."""
    s = make_blank_slide(prs)
    add_text_box(s, *TITLE_BOX, title, pt=TITLE_PT, bold=True)
    # The GIF first frame defines width/height. fit_image_in_box uses PIL.Image
    # which on multi-frame GIFs reports the first frame's size - fine for layout.
    fit_image_in_box(prs, s, gif_path, IMAGE_AREA)
    if caption:
        add_text_box(s, *CAPTION_BOX, caption,
                     pt=CAPTION_PT, color=GREY_RGB, align="center")


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------

def deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ─── Title ──────────────────────────────────────────────────────────────
    add_title_slide(prs)

    # ─── Umbrella overview ─────────────────────────────────────────────────
    umbrella = FIG / "00_umbrella_overview.png"
    if umbrella.exists():
        s = make_blank_slide(prs)
        # Image fills the whole slide with a small border
        fit_image_in_box(prs, s, umbrella,
                         (Inches(0.15), Inches(0.15),
                          Inches(13.033), Inches(7.2)))

    # ─── Stage in/out overview ─────────────────────────────────────────────
    add_section_divider(prs, "Pipeline at a glance",
                        "One slide per stage: real input/output from the traditional run.")

    add_image_slide(
        prs,
        "Stage 1 — stringart: mask → 12 per-angle branches",
        FIG / "02_stage1_branches.png",
        "Left: input mask. Middle: branches color-coded by angle bin. "
        "Right: union of all 12 branches.",
    )
    add_image_slide(
        prs,
        "Stage 2 — preprocess: per-branch cleanup",
        FIG / "03_stage2_clean.png",
        "Top row: a single branch before and after preprocess. "
        "Bottom row: union of all branches before and after.",
    )
    add_image_slide(
        prs,
        "Stage 3 — reconnect: per-angle fragments → bundles",
        FIG / "04_stage3_reconnect.png",
        "Left: cleaned fragments overlaid on SEM. "
        "Right: reconnected bundles (one color per ID).",
    )
    add_image_slide(
        prs,
        "Stage 4 — postprocess: smart-width final bundles",
        FIG / "05_stage4_postprocess.png",
        "Left: stage 3 thin reconnect IDs. "
        "Right: stage 4 rendered at SEM-measured widths (final output).",
    )

    # ─── Algorithm zooms ────────────────────────────────────────────────────
    add_section_divider(prs, "Algorithm zooms",
                        "Geometric primitives the gates and merges operate on.")

    add_image_slide(
        prs,
        "Reconnect gates — what they actually look at",
        FIG / "09_reconnect_gates.png",
        "Three example tip pairs. Yellow crosses = skeleton endpoints. "
        "Gates: dist, forward cosine, opposition, line residual.",
    )
    add_image_slide(
        prs,
        "Smart junction merge (experimental skeleton path)",
        FIG / "07_smart_junction_merge.png",
        "Fuse a collinear arm pair at a 3-arm junction. Anti-parallel tangents (cos < −0.95) qualify.",
    )
    add_image_slide(
        prs,
        "Adaptive binning (experimental skeleton path)",
        FIG / "08_adaptive_binning.png",
        "If a single skeleton piece spans only two adjacent bins, unify all its pixels to the dominant bin.",
    )

    # ─── Stage 1 schematics ─────────────────────────────────────────────────
    add_section_divider(prs, "Stage 1 — stringart, step by step",
                        "Tiled probabilistic Hough with two acceptance gates + multi-grid voting.")
    stage1 = [
        ("Step 1: tile the mask", "01_tile_grid.png",
         "Each tile is processed independently by probabilistic Hough."),
        ("Step 2: Hough candidates inside a tile", "02_hough_in_tile.png",
         "Many candidate segments are drawn; only those passing both gates are kept."),
        ("Gate A: new-pixels added", "03_gate_newpix.png",
         "A candidate must add at least min_accept_newpix previously-unclaimed mask pixels."),
        ("Gate B: support density", "04_gate_density.png",
         "Rejects long fake chords that mostly cross background — must lie on the mask."),
    ]
    for title, fname, caption in stage1:
        add_image_slide(prs, title, SCHEMA / "01_stringart" / fname, caption)

    # Angle-binning step: native-shape version (each bin movable)
    add_native_angle_binning_slide(prs)

    add_image_slide(
        prs,
        "Step 5: multi-grid voting",
        SCHEMA / "01_stringart" / "06_multi_grid_voting.png",
        "Run the stage at N grid offsets; keep only pixels appearing in ≥ vote_min grids.",
    )

    # ─── Stage 2 schematics ─────────────────────────────────────────────────
    add_section_divider(prs, "Stage 2 — preprocess, step by step",
                        "Conservative cleanup of each per-angle branch before reconnect.")
    stage2 = [
        ("Step 1: threshold the grayscale branch", "01_threshold.png",
         "Binarize each per-angle branch with --pre-bin-threshold (default 127)."),
        ("Step 2: drop tiny connected components", "02_drop_specks.png",
         "Components below --pre-min-component pixels are removed (specks, single-pixel noise)."),
        ("Step 3: oriented morphological close", "03_oriented_close.png",
         "A line-shaped kernel along the branch's angle bridges small intra-bundle gaps."),
        ("Step 4: dominant 2-tip path reduction", "04_dominant_path.png",
         "Multi-tip components reduce to their skeleton diameter (the two longest arms)."),
        ("Step 5: B-spline smoothing", "05_bspline.png",
         "Replaces the jagged polyline with a smooth parametric spline before reconnect."),
    ]
    for title, fname, caption in stage2:
        add_image_slide(prs, title, SCHEMA / "02_preprocess" / fname, caption)

    # ─── Stage 3 schematics ─────────────────────────────────────────────────
    add_section_divider(prs, "Stage 3 — reconnect, gate by gate",
                        "Three staged tip-tip merges (clear → strict → relaxed), then smooth passes.")
    stage3 = [
        ("stage_clear — accept", "01_stage_clear_accept.png",
         "Tiny gap, opposition near 1. The cheapest gate, used for obvious continuations."),
        ("stage_clear — reject on distance", "02_stage_clear_reject_distance.png",
         "Geometry perfect but gap exceeds clear_merge_max_dist_px (default 16)."),
        ("stage_strict — accept", "03_stage_strict_accept.png",
         "Moderate gap, all four gates pass (dist, fwd cos, opposition, residual)."),
        ("stage_strict — reject on opposition", "04_stage_strict_reject_opposition.png",
         "Classic T-junction: tips touch, but inward tangents are perpendicular."),
        ("stage_relaxed — accept", "05_stage_relaxed_accept.png",
         "Wider gap, but alignment is still convincing under relaxed thresholds."),
        ("stage_relaxed — reject on line residual", "06_stage_relaxed_reject_residual.png",
         "Two parallel-offset fragments: tangent lines pass far from the gap midpoint."),
        ("Same-layer relaxation", "07_same_layer_relaxation.png",
         "Same-orientation-bin pieces get a looser residual gate."),
        ("Overlap handling: trim vs kill", "08_overlap_trim.png",
         "When two Components overlap heavily, trim keeps the non-overlap fragments; kill drops the smaller."),
        ("Smooth passes — iterate until fixpoint", "09_smooth_pass.png",
         "After staged merges, re-evaluate every neighbour pair; repeat until nothing more accepts."),
    ]
    for title, fname, caption in stage3:
        add_image_slide(prs, title, SCHEMA / "03_reconnect" / fname, caption)

    # ─── Stage 4 schematics ─────────────────────────────────────────────────
    add_section_divider(prs, "Stage 4 — postprocess, step by step",
                        "Re-skeletonize each reconnect layer, smooth, render at its measured SEM width.")
    stage4 = [
        ("Step 1: skeletonize the layer → dominant path", "01_skeletonize_and_path.png",
         "Each multilabel layer is reduced to its endpoint-to-endpoint centerline."),
        ("Step 2: moving-window smoothing", "02_smooth_window.png",
         "The jagged centerline becomes a smooth path (window size = --smooth-window)."),
        ("Step 3: SEM edge sampling along normals", "03_smart_width_sampling.png",
         "At each centerline sample, probe the normal direction for opposite-slope SEM edges."),
        ("Step 4: median width per ID (with outlier trim)", "04_smart_width_median.png",
         "Per-sample widths → median → bundle rendered at that width (fallback: --thicken-px)."),
        ("Step 5: overlap absorb", "05_overlap_absorb.png",
         "Near-duplicate bundles whose intersection ≥ thr × smaller area get merged."),
        ("Step 6: occlusion trim", "06_occlusion_trim.png",
         "Lower-priority layers mostly hidden by earlier layers get trimmed; tiny survivors dropped."),
    ]
    for title, fname, caption in stage4:
        add_image_slide(prs, title, SCHEMA / "04_postprocess" / fname, caption)

    # ─── Bundle tracking section (animated GIFs) ───────────────────────────
    add_section_divider(
        prs, "Following one bundle end-to-end",
        "Each GIF tracks one big bundle through every stage (press F5 to play).",
    )
    add_image_slide(
        prs,
        "Pipeline progression (overall)",
        FIG / "06_pipeline_progression.gif",
        "Mask → SEM → stage 1 → stage 2 → stage 3 → stage 4. "
        "Press F5 to start the slideshow; the GIF will animate.",
    )
    for final_id, dom_branch, rec_id in BUNDLE_GIFS:
        gif = FIG / f"06b_bundle_track_id{final_id}.gif"
        if not gif.exists():
            print(f"  skip missing GIF: {gif.name}")
            continue
        add_gif_slide(
            prs,
            f"Bundle id {final_id}: {dom_branch}, {rec_id}",
            gif,
            "7 frames @ 2.5 s each. The target bundle is highlighted red against the dimmed run "
            "at every stage. Yellow X marks in the stage-3 frame = the skeleton endpoints "
            "reconnect's gates operate on.",
        )

    # ─── Closing slide ──────────────────────────────────────────────────────
    s = make_blank_slide(prs)
    add_text_box(s, Inches(0.5), Inches(2.5), Inches(12.5), Inches(1.5),
                 "Thanks", pt=44, bold=True, color=ACCENT_RGB, align="center")
    add_text_box(s, Inches(0.5), Inches(4.0), Inches(12.5), Inches(2.0),
                 "Markdown source: docs/pipeline_overview.md\n"
                 "Self-contained HTML: docs/pipeline_overview.html\n"
                 "Figure generators: docs/generate_figures.py, docs/generate_schematics.py",
                 pt=16, color=GREY_RGB, align="center")

    prs.save(str(OUT))
    n_slides = len(prs.slides)
    size_mb = OUT.stat().st_size / 1024 / 1024
    print(f"wrote {OUT}  ({n_slides} slides, {size_mb:.1f} MB)")


if __name__ == "__main__":
    deck()
